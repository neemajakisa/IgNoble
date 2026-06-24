#!/usr/bin/env python3
"""scripts/make_slides.py
Generate IgNoble presentation: two slides
  1. The Ig Nobel Prize — what it is, example, problem statement
  2. Updated pipeline diagram
"""

import io
import os
import sys

# Must be set before cairosvg import so cairocffi can find libcairo on macOS
os.environ['DYLD_LIBRARY_PATH'] = '/opt/homebrew/lib:' + os.environ.get('DYLD_LIBRARY_PATH', '')

import cairosvg
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ───────────────────────────────────────────────────────────────────

DARK      = RGBColor(0x1a, 0x17, 0x14)
RED       = RGBColor(0xa8, 0x3b, 0x18)
GREEN     = RGBColor(0x3d, 0x6b, 0x4f)
WARM_BG   = RGBColor(0xf8, 0xf5, 0xf0)
WHITE     = RGBColor(0xff, 0xff, 0xff)
LIGHT_RED = RGBColor(0xfd, 0xf1, 0xee)
LIGHT_GRN = RGBColor(0xeb, 0xf5, 0xee)
WARM_TAN  = RGBColor(0xf0, 0xec, 0xe6)
GOLD      = RGBColor(0xe8, 0xa8, 0x90)
GREY      = RGBColor(0x55, 0x55, 0x55)

PROJ_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SVG_PATH  = os.path.join(PROJ_ROOT, "pipeline.svg")
OUT_PPTX  = os.path.join(PROJ_ROOT, "ig_nobel_slides.pptx")

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None, lw=0.75):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def text_box(slide, x, y, w, h, lines):
    """
    lines = list of dicts with keys:
      text, size, color, bold(opt), italic(opt), align(opt)
    Each dict becomes one paragraph. Use empty text + small size for spacing.
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get('align', PP_ALIGN.LEFT)
        run = p.add_run()
        run.text           = spec.get('text', '')
        run.font.size      = Pt(spec.get('size', 11))
        run.font.color.rgb = spec.get('color', DARK)
        run.font.bold      = spec.get('bold', False)
        run.font.italic    = spec.get('italic', False)
    return box


# ── Convert SVG to PNG ────────────────────────────────────────────────────────

with open(SVG_PATH, 'rb') as f:
    svg_bytes = f.read()

png_bytes = cairosvg.svg2png(bytestring=svg_bytes, scale=2.0)
print("✓ SVG → PNG")

# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

W = prs.slide_width
H = prs.slide_height

# ─── SLIDE 1: The Ig Nobel Prize ─────────────────────────────────────────────

s1 = prs.slides.add_slide(blank)
s1.background.fill.solid()
s1.background.fill.fore_color.rgb = WARM_BG

BAR_H = Inches(1.12)

# Header bar
add_rect(s1, 0, 0, W, BAR_H, DARK)
add_rect(s1, 0, BAR_H, W, Inches(0.04), RED)  # red accent line

text_box(s1, Inches(0.35), Inches(0.06), Inches(12.5), Inches(0.62), [
    {'text': 'THE IG NOBEL PRIZE', 'size': 30, 'color': WHITE, 'bold': True},
])
text_box(s1, Inches(0.35), Inches(0.67), Inches(12.5), Inches(0.38), [
    {'text': '"First make people LAUGH, then make them THINK"',
     'size': 13, 'color': GOLD, 'italic': True},
])

# Three columns
COL_W = Inches(3.95)
GAP   = Inches(0.22)
X1    = Inches(0.27)
X2    = X1 + COL_W + GAP
X3    = X2 + COL_W + GAP
TOP   = BAR_H + Inches(0.09)
CH    = H - TOP - Inches(0.05)
HDR_H = Inches(0.40)
PAD   = Inches(0.14)

add_rect(s1, X1, TOP, COL_W, CH, WARM_TAN)
add_rect(s1, X2, TOP, COL_W, CH, LIGHT_RED)
add_rect(s1, X3, TOP, COL_W, CH, LIGHT_GRN)

add_rect(s1, X1, TOP, COL_W, HDR_H, DARK)
add_rect(s1, X2, TOP, COL_W, HDR_H, RED)
add_rect(s1, X3, TOP, COL_W, HDR_H, GREEN)

S = 11   # body font size
G = 3.5  # gap paragraph font size (creates vertical space)

for x, label in ((X1, 'WHAT IS IT?'), (X2, 'CLASSIC EXAMPLE'), (X3, 'OUR PROBLEM STATEMENT')):
    text_box(s1, x + PAD, TOP + Inches(0.06), COL_W - PAD*2, HDR_H, [
        {'text': label, 'size': 10.5, 'color': WHITE, 'bold': True},
    ])

CY = TOP + HDR_H + Inches(0.13)
CW = COL_W - PAD * 2
CXH = CH - HDR_H - Inches(0.13)

# Column 1 — What is it?
text_box(s1, X1 + PAD, CY, CW, CXH, [
    {'text': 'Annual prizes for genuine scientific achievements that are simultaneously funny AND real.',
     'size': S, 'color': DARK},
    {'text': '', 'size': G, 'color': DARK},
    {'text': '•  Real research, real journals, real scientists', 'size': S, 'color': DARK},
    {'text': '•  Administered by the Annals of Improbable Research since 1991', 'size': S, 'color': DARK},
    {'text': '•  Awarded at Harvard in an actual ceremony', 'size': S, 'color': DARK},
    {'text': '•  Categories: Medicine, Physics, Chemistry, Biology, Peace, Economics…',
     'size': S, 'color': DARK},
    {'text': '', 'size': G, 'color': DARK},
    {'text': 'The humor comes FROM the scientific finding — not from how it is described.',
     'size': S, 'color': RED, 'bold': True},
    {'text': '', 'size': G, 'color': DARK},
    {'text': 'The laugh must arrive BEFORE the explanation does.',
     'size': S, 'color': DARK, 'italic': True},
])

# Column 2 — Classic Example
text_box(s1, X2 + PAD, CY, CW, CXH, [
    {'text': '2000 Ig Nobel Prize in Physics', 'size': S, 'color': RED, 'bold': True},
    {'text': '', 'size': G, 'color': DARK},
    {'text': '“On the Possibility of a Frog to Be Levitated by a Magnet”',
     'size': S, 'color': DARK, 'italic': True},
    {'text': '', 'size': G, 'color': DARK},
    {'text': 'André Geim used high-field magnets to levitate a living frog against gravity.',
     'size': S, 'color': DARK},
    {'text': '', 'size': G, 'color': DARK},
    {
        'text': '→  Geim later won the 2010 Nobel Prize in Physics (for graphene!)',
        'size': S, 'color': GREEN, 'bold': True,
    },
    {'text': '', 'size': G + 1, 'color': DARK},
    {'text': 'Other famous winners:', 'size': S, 'color': DARK, 'bold': True},
    {'text': '•  Slipping on banana peels (Physics, 2014)', 'size': S, 'color': DARK},
    {'text': '•  Armadillos as leprosy transmitters (Medicine, 2012)', 'size': S, 'color': DARK},
    {'text': '•  Why woodpeckers don’t get headaches (Medicine, 2006)', 'size': S, 'color': DARK},
])

# Column 3 — Problem Statement
text_box(s1, X3 + PAD, CY, CW, CXH, [
    {'text': 'Can a multi-agent AI pipeline autonomously generate and evaluate Ig Nobel-caliber research proposals?',
     'size': S, 'color': DARK, 'bold': True},
    {'text': '', 'size': G, 'color': DARK},
    {'text': 'Challenges:', 'size': S, 'color': GREEN, 'bold': True},
    {'text': '', 'size': G - 1, 'color': DARK},
    {'text': '•  The “laugh then think” criterion is subtle — hard to encode in a prompt',
     'size': S, 'color': DARK},
    {'text': '•  Requires novelty + scientific plausibility + inherent absurdity simultaneously',
     'size': S, 'color': DARK},
    {'text': '•  Distinguishing genuinely funny from merely taboo, niche, or weird',
     'size': S, 'color': DARK},
    {'text': '•  Novelty detection requires real-time web search', 'size': S, 'color': DARK},
    {'text': '', 'size': G, 'color': DARK},
    {'text': 'Our approach:', 'size': S, 'color': GREEN, 'bold': True},
    {'text': '', 'size': G - 1, 'color': DARK},
    {'text': '4 specialised agents with distinct roles, web search for novelty checking, and iterative critique–revision loops.',
     'size': S, 'color': DARK},
])


# ─── SLIDE 2: Pipeline ────────────────────────────────────────────────────────

s2 = prs.slides.add_slide(blank)
s2.background.fill.solid()
s2.background.fill.fore_color.rgb = WARM_BG

add_rect(s2, 0, 0, W, Inches(0.68), DARK)
add_rect(s2, 0, Inches(0.68), W, Inches(0.035), RED)
text_box(s2, Inches(0.35), Inches(0.1), Inches(12.5), Inches(0.54), [
    {'text': 'Our Multi-Agent Pipeline', 'size': 22, 'color': WHITE, 'bold': True},
])

# Fit pipeline image below header, centred
img_top = Inches(0.75)
avail_w = W - Inches(0.2)
avail_h = H - img_top - Inches(0.05)

svg_w, svg_h = 1000, 760
if svg_w / svg_h > avail_w / avail_h:
    img_w = avail_w
    img_h = img_w * svg_h / svg_w
else:
    img_h = avail_h
    img_w = img_h * svg_w / svg_h

img_x = (W - img_w) / 2
img_y = img_top + (avail_h - img_h) / 2

s2.shapes.add_picture(io.BytesIO(png_bytes), img_x, img_y, img_w, img_h)

# ── Save ──────────────────────────────────────────────────────────────────────

prs.save(OUT_PPTX)
print(f"✓ Saved  {OUT_PPTX}")
